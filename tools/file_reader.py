"""Multi-format file reader tools for the intake agent.

Supported: .pdf, .docx, .txt, .md, .pptx, .xlsx, .xls, images.

PDF reading strategy (vision-first):
  read_pdf_smart  — renders each page → PNG → vision LLM (gpt-5.4-mini).
                    Handles text, tables, flowcharts, diagrams accurately.
  read_pdf        — pypdf text-only fallback (fast, no API cost).

Image reading: describe_image / describe_requirement_image use vision LLM
with a requirements-analyst prompt to extract structured content.
"""
from __future__ import annotations

import base64
import os
from pathlib import Path
from typing import Optional

from langchain_core.tools import tool
from loguru import logger as _log

_fr_log = _log.bind(ctx="file_reader")

# ── Vision prompt shared across PDF pages and standalone images ───────────────

_VISION_PROMPT = """\
You are a senior business analyst extracting requirements from a document page.
Extract ALL content visible on this page completely and accurately.
Focus on:
- Problem statement / pain points (current state that needs solving)
- Business objective (what the solution should achieve)
- Functional requirements (what the system must DO — list each explicitly)
- Non-functional requirements (performance targets, security, availability — \
include ALL numbers and units)
- Data sources and integrations (systems, APIs, databases mentioned)
- Expected outcomes and success metrics (preserve all numbers and units exactly)
- Stakeholders and user roles
- Constraints, assumptions, timelines, budget hints
- ALL table data — reproduce as Markdown tables
- Diagrams / flowcharts / architecture — describe entities, flows, relationships
Output as structured Markdown. Preserve every number, unit, and technical term exactly."""


# ── Helpers ───────────────────────────────────────────────────────────────────

def _size_guard(path: Path, max_mb: float = 100.0) -> None:
    mb = path.stat().st_size / 1_048_576
    if mb > max_mb:
        raise ValueError(f"File too large: {mb:.1f} MB (max {max_mb} MB)")


def _resolve(path_str: str) -> Path:
    from .workspace import resolve_path, get_workspace, get_input_dir, _thread_input_registry, _active_thread_id_var
    p = resolve_path(path_str).expanduser()
    if p.exists():
        return p
    # Fallback: check workspace directory for the same filename
    # (handles cases where render saved to workspace but agent points to output)
    ws_copy = get_workspace() / p.name
    if ws_copy.exists():
        _fr_log.warning(f"_resolve: {path_str!r} not found, using workspace copy {ws_copy}")
        return ws_copy
    # Fallback: if path is absolute and under a known input dir, try it directly
    abs_path = Path(path_str)
    if abs_path.is_absolute() and abs_path.exists():
        return abs_path
    # Fallback: search for filename in the thread's input directory
    fname = Path(path_str).name
    input_dir = get_input_dir()
    if input_dir is None:
        tid = _active_thread_id_var.get()
        if tid:
            input_dir = _thread_input_registry.get(tid)
    if input_dir:
        candidate = input_dir / fname
        if candidate.exists():
            _fr_log.warning(f"_resolve: found {fname!r} by name in input_dir {input_dir}")
            return candidate
    raise FileNotFoundError(f"File not found: {path_str}")


def _b64_png(img_bytes: bytes) -> str:
    return base64.b64encode(img_bytes).decode()


def _vision_call(b64_image: str, prompt: str, model: str = "gpt-5.4-mini") -> str:
    from packages.config import get_llm
    from langchain_openai import ChatOpenAI
    llm = get_llm("intake_agent")
    if isinstance(llm, str):
        llm = ChatOpenAI(model=model)
    response = llm.invoke([{
        "role": "user",
        "content": [
            {"type": "image_url",
             "image_url": {"url": f"data:image/png;base64,{b64_image}"}},
            {"type": "text", "text": prompt},
        ],
    }])
    return str(response.content)


# ── Tools ─────────────────────────────────────────────────────────────────────

@tool
def list_input_files(directory: str) -> str:
    """List all supported requirement files in a directory (recursive).

    Returns a formatted list of file paths and sizes.
    """
    from .workspace import resolve_path, get_input_single_file
    SUPPORTED = {".pdf", ".docx", ".txt", ".md", ".pptx", ".xlsx", ".xls",
                 ".png", ".jpg", ".jpeg", ".webp"}
    d = resolve_path(directory).expanduser()
    single_file = get_input_single_file()
    _fr_log.info(f"list_input_files | {directory!r} → resolved={d}"
                 + (f"  [single={single_file}]" if single_file else ""))
    if not d.exists():
        _fr_log.warning(f"  Directory not found: {d}")
        return f"Directory not found: {directory}"
    lines = []
    for p in sorted(d.rglob("*")):
        if not (p.suffix.lower() in SUPPORTED and p.is_file()):
            continue
        # If a single-file filter is set, only return that file
        if single_file and p.name != single_file:
            continue
        kb = p.stat().st_size / 1024
        # Return virtual /input/<relative> path so the LLM uses the correct
        # virtual prefix when calling read tools (not the raw absolute path).
        try:
            rel = p.relative_to(d)
            virtual = f"/input/{rel}"
        except ValueError:
            virtual = f"/input/{p.name}"
        lines.append(f"{virtual}  ({kb:.0f} KB)")
    _fr_log.info(f"  Found {len(lines)} file(s)")
    return "\n".join(lines) if lines else "No supported files found."


@tool
def read_pdf_smart(
    file_path: str,
    vision_model: str = "gpt-5.4-mini",
    dpi: int = 150,
) -> str:
    """Extract ALL content from a PDF using vision AI per page.

    Renders each page to PNG and passes it to a vision LLM. Accurately captures
    text, tables, diagrams, flowcharts, and annotations — far better than
    text-only extraction for complex requirement documents and business cases.

    Args:
        file_path:    Absolute path to the PDF file.
        vision_model: Vision-capable model to use (default gpt-5.4-mini).
        dpi:          Render resolution — 150 is balanced; use 200 for dense text.
    """
    p = _resolve(file_path)
    _size_guard(p, max_mb=100)
    _fr_log.info(f"read_pdf_smart | {p.name}  model={vision_model}  dpi={dpi}")
    try:
        import fitz  # PyMuPDF
    except ImportError:
        _fr_log.warning("PyMuPDF not found — falling back to text-only read_pdf")
        return read_pdf.invoke({"file_path": file_path})

    try:
        doc = fitz.open(str(p))
        total = len(doc)
        _fr_log.info(f"  {total} pages, processing via vision LLM…")
        scale = dpi / 72.0
        mat = fitz.Matrix(scale, scale)
        parts = []

        for page_num in range(total):
            page = doc[page_num]
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            b64 = _b64_png(img_bytes)
            page_content = _vision_call(b64, _VISION_PROMPT, model=vision_model)
            parts.append(f"## [Page {page_num + 1} / {total}]\n\n{page_content}")
            _fr_log.debug(f"  page {page_num+1}/{total} done ({len(page_content)} chars)")

        doc.close()
        _fr_log.success(f"  read_pdf_smart done — {total} pages, {sum(len(p) for p in parts):,} chars")
        return f"# {p.name}\n\n" + "\n\n---\n\n".join(parts)

    except Exception as e:
        _fr_log.error(f"  PDF smart extraction error: {e}")
        return f"[PDF smart extraction error: {e}]"


@tool
def read_pdf(file_path: str, max_chars: int = 200_000) -> str:
    """Extract text from a PDF file (text-only, no vision).

    Fast fallback when vision is not needed (plain text PDFs).
    For complex documents with tables/diagrams, use read_pdf_smart instead.
    """
    p = _resolve(file_path)
    _size_guard(p)
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(p))
        parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            parts.append(f"[Page {i+1}]\n{text}")
        return "\n\n".join(parts)[:max_chars]
    except Exception as e:
        return f"[PDF extraction error: {e}]"


@tool
def read_docx(file_path: str) -> str:
    """Extract text and tables from a Word (.docx) file, preserving headings."""
    p = _resolve(file_path)
    _size_guard(p, max_mb=20)
    _fr_log.info(f"read_docx | {p.name}")
    try:
        from docx import Document
        doc = Document(str(p))
        lines = []
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            if "Heading 1" in style:
                prefix = "## "
            elif "Heading" in style:
                prefix = "### "
            else:
                prefix = ""
            if para.text.strip():
                lines.append(f"{prefix}{para.text.strip()}")

        for table in doc.tables:
            if not table.rows:
                continue
            rows = []
            for i, row in enumerate(table.rows):
                cells = [c.text.strip() for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows.append("| " + " | ".join("---" for _ in cells) + " |")
            lines.append("\n".join(rows))

        result = "\n\n".join(lines)
        _fr_log.success(f"  read_docx done — {len(result):,} chars")
        return result
    except Exception as e:
        _fr_log.error(f"  DOCX extraction error: {e}")
        return f"[DOCX extraction error: {e}]"


@tool
def read_txt(file_path: str, max_chars: int = 100_000) -> str:
    """Read a plain text or Markdown file."""
    p = _resolve(file_path)
    _size_guard(p, max_mb=20)
    return p.read_text(encoding="utf-8", errors="replace")[:max_chars]


@tool
def read_pptx(file_path: str) -> str:
    """Extract text from PowerPoint slides (title + body + speaker notes)."""
    p = _resolve(file_path)
    _size_guard(p)
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"Notes: {notes}")
            slides.append("\n".join(parts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


@tool
def read_xlsx(file_path: str, max_rows: int = 500) -> str:
    """Extract data from an Excel file as Markdown tables (up to 500 rows/sheet)."""
    p = _resolve(file_path)
    _size_guard(p)
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data: list[list[str]] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    break
                rows_data.append(["" if v is None else str(v) for v in row])

            if not rows_data:
                continue

            # Format as Markdown table
            md_rows = []
            for i, row in enumerate(rows_data):
                md_rows.append("| " + " | ".join(row) + " |")
                if i == 0:
                    md_rows.append("| " + " | ".join("---" for _ in row) + " |")
            if len(rows_data) >= max_rows:
                md_rows.append(f"_(truncated at {max_rows} rows)_")

            parts.append(f"### Sheet: {sheet_name}\n\n" + "\n".join(md_rows))
        return "\n\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


@tool
def describe_image(
    file_path: str,
    context: str = "",
    model: str = "gpt-5.4-mini",
) -> str:
    """Extract requirements, diagrams, and data from an image file using vision AI.

    Args:
        file_path: Absolute path to PNG, JPG, or WEBP image.
        context:   Optional context hint (e.g. "architecture diagram for voicebot").
        model:     Vision model to use.
    """
    p = _resolve(file_path)
    mb = p.stat().st_size / 1_048_576
    if mb > 20:
        return f"[Image too large: {mb:.1f} MB (max 20 MB)]"

    try:
        with open(p, "rb") as f:
            raw = f.read()

        suffix = p.suffix.lower().lstrip(".")
        media_type = {
            "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png", "webp": "image/webp",
        }.get(suffix, "image/png")

        b64 = base64.b64encode(raw).decode()

        prompt = _VISION_PROMPT
        if context:
            prompt = f"Context: {context}\n\n{prompt}"

        from packages.config import get_llm
        from langchain_openai import ChatOpenAI
        llm = get_llm("intake_agent")
        if isinstance(llm, str):
            llm = ChatOpenAI(model=llm)
        response = llm.invoke([{
            "role": "user",
            "content": [
                {"type": "image_url",
                 "image_url": {"url": f"data:{media_type};base64,{b64}"}},
                {"type": "text", "text": prompt},
            ],
        }])
        return str(response.content)
    except Exception as e:
        return f"[Image description error: {e}]"


@tool
def write_raw_features(content: str) -> str:
    """Save the structured requirement analysis to /workspace/raw_features.md.

    Call this ONCE after reading and analyzing ALL input files.
    The content must follow the 11-section schema defined in the intake prompt.
    """
    from .workspace import write_text, RAW_FEATURES_FILE, get_workspace
    _fr_log.info(f"write_raw_features | {len(content):,} chars → workspace={get_workspace()}")
    write_text(RAW_FEATURES_FILE, content)
    _fr_log.success(f"  raw_features.md saved ✓")
    return f"Saved {len(content):,} chars to raw_features.md"
